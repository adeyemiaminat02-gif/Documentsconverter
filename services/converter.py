import os
import fitz  # PyMuPDF
import pdfplumber
from docx import Document
import pypandoc
import markdown
from bs4 import BeautifulSoup
from openpyxl import Workbook, load_workbook
from pptx import Presentation
import weasyprint
from utils.logger import logger

VALID_CONVERSIONS = {
    'PDF': ['DOCX', 'TXT', 'HTML', 'MD'],
    'DOCX': ['PDF', 'TXT', 'HTML', 'MD'],
    'TXT': ['PDF', 'DOCX', 'HTML'],
    'HTML': ['PDF', 'DOCX', 'MD'],
    'MD': ['HTML', 'PDF', 'DOCX'],
    'EPUB': ['PDF', 'DOCX', 'TXT'],
    'RTF': ['PDF'],
    'ODT': ['DOCX'],
    'CSV': ['XLSX'],
    'XLSX': ['CSV'],
    'PPTX': ['PDF']
}

def get_valid_targets(src_ext: str) -> list:
    return VALID_CONVERSIONS.get(src_ext.upper().replace('.', ''), [])

def convert_file(src_path: str, target_format: str, output_dir: str) -> str:
    base_name = os.path.splitext(os.path.basename(src_path))[0]
    src_ext = os.path.splitext(src_path)[1].upper().replace('.', '')
    target_format = target_format.upper()
    dest_path = os.path.join(output_dir, f"{base_name}.{target_format}.all")
    dest_path = dest_path.replace(".all", "") 

    logger.info(f"Converting {src_path} ({src_ext}) to {target_format}")

    # 1. PDF Sub-Engine
    if src_ext == 'PDF':
        if target_format == 'TXT':
            with pdfplumber.open(src_path) as pdf:
                text = "".join([page.extract_text() or "" for page in pdf.pages])
            with open(dest_path, "w", encoding="utf-8") as f:
                f.write(text)
            return dest_path
        elif target_format == 'DOCX':
            doc = Document()
            with pdfplumber.open(src_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        doc.add_paragraph(text)
            doc.save(dest_path)
            return dest_path
        elif target_format in ['HTML', 'MD']:
            text_content = ""
            with pdfplumber.open(src_path) as pdf:
                text_content = "".join([page.extract_text() or "" for page in pdf.pages])
            if target_format == 'HTML':
                html_data = f"<html><body><pre>{text_content}</pre></body></html>"
                with open(dest_path, "w", encoding="utf-8") as f:
                    f.write(html_data)
            else:
                with open(dest_path, "w", encoding="utf-8") as f:
                    f.write(text_content)
            return dest_path

    # 2. DOCX Sub-Engine
    if src_ext == 'DOCX':
        if target_format == 'PDF':
            # Uses Weasyprint via temporary HTML abstraction layer to guard formatting layout
            doc = Document(src_path)
            html_content = "".join([f"<p>{p.text}</p>" for p in doc.paragraphs])
            weasyprint.HTML(string=html_content).write_pdf(dest_path)
            return dest_path
        elif target_format == 'TXT':
            doc = Document(src_path)
            text = "\n".join([p.text for p in doc.paragraphs])
            with open(dest_path, "w", encoding="utf-8") as f:
                f.write(text)
            return dest_path
        elif target_format in ['HTML', 'MD']:
            doc = Document(src_path)
            text = "\n".join([p.text for p in doc.paragraphs])
            if target_format == 'HTML':
                content = "".join([f"<p>{p.text}</p>" for p in doc.paragraphs])
                with open(dest_path, "w", encoding="utf-8") as f:
                    f.write(f"<html><body>{content}</body></html>")
            else:
                with open(dest_path, "w", encoding="utf-8") as f:
                    f.write(text)
            return dest_path

    # 3. Plain Text Sub-Engine
    if src_ext == 'TXT':
        with open(src_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
        if target_format == 'PDF':
            html_content = f"<html><body><pre>{content}</pre></body></html>"
            weasyprint.HTML(string=html_content).write_pdf(dest_path)
            return dest_path
        elif target_format == 'DOCX':
            doc = Document()
            doc.add_paragraph(content)
            doc.save(dest_path)
            return dest_path
        elif target_format == 'HTML':
            with open(dest_path, "w", encoding="utf-8") as f:
                f.write(f"<html><body><pre>{content}</pre></body></html>")
            return dest_path

    # 4. Markup Sub-Engine (HTML / Markdown / EPUB)
    if src_ext in ['HTML', 'MD', 'EPUB', 'RTF', 'ODT']:
        try:
            # Fallback wrapper via internal pypandoc execution pipeline wrappers
            pypandoc.convert_file(src_path, target_format.lower(), outputfile=dest_path)
            return dest_path
        except Exception as e:
            logger.warning(f"Pandoc lifecycle step missed/failed fallback to parsing core. Engine error: {e}")
            if src_ext == 'MD' and target_format == 'HTML':
                with open(src_path, "r", encoding="utf-8") as f:
                    html_text = markdown.markdown(f.read())
                with open(dest_path, "w", encoding="utf-8") as f:
                    f.write(html_text)
                return dest_path
            raise RuntimeError("Engine processing constraint or format binary missing.")

    # 5. Data Struct-Sheet Layout Sub-Engine
    if src_ext == 'CSV' and target_format == 'XLSX':
        wb = Workbook()
        ws = wb.active
        import csv
        with open(src_path, "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            for row in reader:
                ws.append(row)
        wb.save(dest_path)
        return dest_path

    if src_ext == 'XLSX' and target_format == 'CSV':
        wb = load_workbook(src_path, data_only=True)
        ws = wb.active
        import csv
        with open(dest_path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            for row in ws.iter_rows(values_only=True):
                writer.append(list(row) if row else [])
        return dest_path

    # 6. Presentation Layout Slideware Sub-Engine
    if src_ext == 'PPTX' and target_format == 'PDF':
        prs = Presentation(src_path)
        slides_text = []
        for slide in prs.slides:
            text_chunks = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_chunks.append(shape.text)
            slides_text.append("<br>".join(text_chunks))
        html_content = "".join([f"<div style='page-break-after:always;'>{st}</div>" for st in slides_text])
        weasyprint.HTML(string=html_content).write_pdf(dest_path)
        return dest_path

    raise ValueError(f"Unsupported conversion matrix: {src_ext} to {target_format}")
